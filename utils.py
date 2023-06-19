import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns
import pandas as pd
from datetime import datetime
from  matplotlib.colors import LinearSegmentedColormap
import dataframe_image as dfi

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from io import BytesIO

dict_langs2 = {'es': 'spanish',
             'pt': 'portuguese',
             'en': 'english',
             'id': 'indonesian',
             'he': 'israel',
             'ms': 'english'}


def transform_inches(x):
  return x / 2.56

def create_cmap(*colors):
  lis_colors = list(colors)
  cmap = LinearSegmentedColormap.from_list('gr',lis_colors, N=256)
  return cmap, cmap.reversed()

def transformation(dataframe):
  df = dataframe.copy()
  df.columns = df.columns.str.lower().str.replace(' ', '_').str.replace("'",'')
  df['assignee'] = df['reviewers_email'].str.split("@").str[0]
  df = df.dropna(subset=["reviewers_email"]).reset_index(drop=True)
  for col in ['upload_date', 'review_start', 'review_finish']:
      df[col] = pd.to_datetime(df[col],format="%d.%m.%Y %H:%M:%S")
  df["time_reviewing_attempt"] = df['review_finish'] - df['upload_date']
  df["secs_reviewing_attempt"] = df['time_reviewing_attempt'].dt.total_seconds()
  df['sla_issue'] = np.where(df['secs_reviewing_attempt'] > 86400, 1, 0)

  df['date'] = pd.to_datetime(df['review_finish'].dt.date)
  df['week'] = df["date"].dt.to_period('W').dt.start_time
  df['ticket_id'] = df['ticket_id'].astype('int64')
  df['lang_reg'] = df['language'].str.lower()
  df['language'] = df['language'].str.extract(pat="(\w+)(?: \- \w+)")[0].str.lower()
  df['language'] = df['language'].map(dict_langs2)
  return df.reset_index(drop=True)

def filter_days(df, col="date", num_days=7):
  filt = df[df[col] >= (datetime.today() - pd.to_timedelta(num_days, "D"))]
  return filt.reset_index(drop=True)



def heatmap_function(df, transformed=False, num_days=7):
  if not transformed:
    data = transformation(df).copy()
    filt = filter_days(data, num_days=num_days).copy()
  else:
    filt = df.copy()
  fig, ax = plt.subplots(5, 4, figsize=(transform_inches(24.14), transform_inches(5.68)))
  langs = data['language'].unique()
  for i in range(5):
    for j in range(4):
      filtro = filt[filt['language'] == langs[j]]
      filtro = (filtro.groupby('date').agg({'secs_reviewing_attempt':'mean',
                          'ticket_id':'count',
                          'sla_issue':'sum'}))
      filtro['time'] = np.ceil(filtro['secs_reviewing_attempt'] / 3600)
      filtro.drop(columns='secs_reviewing_attempt', inplace=True)
      filtro['perc_sla'] = filtro['sla_issue'] / filtro['ticket_id']
      filtro.columns = ['revs_made', 'sla_issues', 'avg_time_per_review (hrs)', 'sla_perc_revs_made']
      filtro = filtro[['avg_time_per_review (hrs)','revs_made', 'sla_issues','sla_perc_revs_made']]
      lista_cols = []
      for column in filtro.columns:
        lista_cols.append(pd.DataFrame(filtro[column]).T)

      cbar=True
      if j != 3:
        cbar = False

      if i == 0:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='.0f',  vmin=0, vmax=24, cmap=cmap, cbar=cbar)
        ax[i, j].set_title(langs[j])
      elif i == 1:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='.0f',  vmin=0, vmax=8, cmap=cmap2, cbar=cbar)
      elif i == 2:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='.0f',  vmin=0, vmax=2, cmap=cmap, cbar=cbar)
      else:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='0.1%',  vmin=0, vmax=0.2, cmap=cmap, cbar=cbar)

      ticks_x = []
      if i != 3:
        ax[i, j].set_xticks([])
      else:
        ax[i, j].set_xticklabels(pd.to_datetime(lista_cols[i].columns).strftime('%d-%b'))

      if j != 0:
        ax[i, j].set_yticks([])

      else:
        ax[i, j].tick_params(axis='y', rotation=0, left=False)
      ax[i, j].set_xlabel('')
  fig.suptitle("Reviewers performance", size=16)


  # cmap red, green, yellow

  plt.tight_layout()
  #plt.subplots_adjust(wspace=0.05, hspace=0.2)

  fig.savefig("heatmap.png", bbox_inches='tight')
  plt.show()

def barplot_function(df, transformed=False, num_days=7):
  if not transformed:
    data = transformation(df).copy()
    filt = filter_days(data, num_days=num_days).copy()
  else:
    filt = df.copy()

  fig, ax = plt.subplots(1, 1, figsize=(10.75, 8.03))
  pivot = filt.pivot_table(index='date', columns='language', values='ticket_id', aggfunc='count')
  pivot_2 = filt.pivot_table(index='date', values='assignee', aggfunc='nunique')
  pivot.plot(kind='bar', stacked=True, cmap=cmap2, ax=ax)
  ax.set_xticklabels(pd.to_datetime(pivot.index).strftime('%d-%b'))
  ax.spines['top'].set_visible(False)
  ax.spines['right'].set_visible(False)
  ax.set_ylabel('# projects')
  ax.set_title("Number of projects solved in last 7 days", size=20)

  contador = 0
  for p in ax.patches:
      width, height = p.get_width(), p.get_height()
      x, y = p.get_xy()
      if height != 0:
        height = int(height)
        if contador > 17:
          ax.text(x + width / 2, y + height / 2, height, ha='center', va='center', color='white', size=14)
        else:
          ax.text(x + width / 2, y + height / 2, height, ha='center', va='center', size=14)
      contador += 1

  plt.show()


  fig.savefig('bar_plot.png', bbox_inches='tight')


def top_reviewers(df, transformed=False, col_group='language', num_days=7):
  if not transformed:
    data = transformation(df)
    filt = filter_days(data, num_days=num_days).copy()
  else:
    filt = df.copy()

  filt = filt[filt['attempt_status'] == 'approved']
  top = (filt
         .groupby([col_group, 'assignee'])['ticket_id']
         .count()
         .reset_index()
         .sort_values(by=[col_group, 'ticket_id'], ascending=False))

  top['pos'] = np.nan

  df_sep = []
  for cont in top[col_group].unique():
    top.loc[top[col_group] == cont, 'pos'] = range(1, len(top[top[col_group] == cont]) + 1)
    df_sep.append((top
           .loc[top[col_group] == cont]
           .drop(columns=col_group)
           .set_index('pos'))
           #.head(3))
           .rename(columns={'ticket_id':'nums'})
           .fillna(''))

  for idx, dataf in enumerate(df_sep):
    dataf.index = dataf.index.astype('int')
    dataf.columns = [top[col_group].unique()[idx], 'nums ']


  res = pd.concat(df_sep, axis=1).fillna('')
  dfi.export(res.replace('',0).astype({'nums ':'int'}).replace(0, ' '), 
             f'table_{col_group}.png', 
             table_conversion="matplotlib")
  #testing = filt.groupby(['assignee']).agg({'ticket_id':'count'}).sort_values(by='ticket_id', ascending=False).reset_index().head(5).rename(columns={'project':'nums'})
  #testing.index = list(range(1, 6))
  #dfi.export(testing, f'top_reviewers.png', table_conversion='matplotlib')

  return res.replace('',0).astype({'nums ':'int'}).replace(0, ' ')

def heatmap_function(df, transformed=False, num_days=7):
  if not transformed:
    data = transformation(df).copy()
    filt = filter_days(data, num_days=num_days).copy()
  else:
    filt = df.copy()

  
  langs = data['language'].unique()
  len_langs = len(langs) if len(langs) != 1 else 2
  fig, ax = plt.subplots(4, len(langs), figsize=(24.14, 5.68))
  for i in range(4):
    for j in range(len(langs)):
      filtro = filt[filt['language'] == langs[j]]
      filtro = (filtro.groupby('date').agg({'secs_reviewing_attempt':'mean',
                          'ticket_id':'count',
                          'sla_issue':'sum'}))
      filtro['time'] = np.ceil(filtro['secs_reviewing_attempt'] / 3600)
      filtro.drop(columns='secs_reviewing_attempt', inplace=True)
      filtro['perc_sla'] = filtro['sla_issue'] / filtro['ticket_id']
      filtro.columns = ['revs_made', 'sla_issues', 'avg_time_per_review (hrs)', 'sla_perc_revs_made']
      filtro = filtro[['avg_time_per_review (hrs)','revs_made', 'sla_issues','sla_perc_revs_made']]
      lista_cols = []
      for column in filtro.columns:
        lista_cols.append(pd.DataFrame(filtro[column]).T)

      cbar = True
      if j != len(langs) - 1:
        cbar = False

      if i == 0:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='.0f',  vmin=0, vmax=24, cmap=cmap, cbar=cbar, annot_kws={'fontsize': 15})
        ax[i, j].set_title(langs[j], size=14)
      elif i == 1:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='.0f',  vmin=0, vmax=8, cmap=cmap2, cbar=cbar, annot_kws={'fontsize': 15})
      elif i == 2:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='.0f',  vmin=0, vmax=2, cmap=cmap, cbar=cbar, annot_kws={'fontsize': 15})
      else:
        sns.heatmap(lista_cols[i], ax=ax[i, j], annot=True, fmt='0.1%',  vmin=0, vmax=0.2, cmap=cmap, cbar=cbar, annot_kws={'fontsize': 15})

      ticks_x = []
      if i != 3:
        ax[i, j].set_xticks([])
      else:
        ax[i, j].set_xticklabels(pd.to_datetime(lista_cols[i].columns).strftime('%d-%b'))

      if j != 0:
        ax[i, j].set_yticks([])

      else:
        ax[i, j].tick_params(axis='y', rotation=0, left=False, labelsize=14)
      ax[i, j].set_xlabel('')
  fig.suptitle("Reviewers performance", size=20)


  # cmap red, green, yellow

  plt.tight_layout()
  #plt.subplots_adjust(wspace=0.05, hspace=0.2)

  fig.savefig("heatmap.png", bbox_inches='tight')
  plt.show()

def number_reviewers(df):
  data = transformation(df.copy())
  filt = filter_days(data)
  return filt['assignee'].nunique()

def report(df):
  df = df.dropna()
  tabla = top_reviewers(df, transformed=False, col_group='language', num_days=7)
  barplot_function(df, transformed=False, num_days=7)
  heatmap_function(df, transformed=False, num_days=7)
  filas_tabla = tabla.shape[0]


  prs = Presentation()
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)

  # title
  left = Inches(transform_inches(0.65))
  top = Inches(transform_inches(0.65))
  width = Inches(transform_inches(10.75))
  height = Inches(transform_inches(1.58))
  txBox = slide.shapes.add_textbox(left, top, width, height)
  tf = txBox.text_frame
  p = tf.add_paragraph()
  p.text = f"Report: {datetime.today().strftime('%Y-%m-%d')} - NM"
  p.font.size = Pt(25)
  p.font.bold = True

  #subtitle
  left = Inches(transform_inches(0.65))
  top = Inches(transform_inches(1.65))
  width = Inches(transform_inches(9.5))
  height = Inches(transform_inches(0.9))
  txBox = slide.shapes.add_textbox(left, top, width, height)
  tf = txBox.text_frame
  p = tf.add_paragraph()
  p.text = "Information about code reviewers performance."
  p.font.size = Pt(9)
  p.font.bold = True

  #shapes

  left = Inches(transform_inches(10.75 + 0.65 + 0.65))
  top = Inches(transform_inches(0.65))
  width = Inches(transform_inches(6.11))
  height = Inches(transform_inches(2.11))

  shape = slide.shapes.add_shape(
      MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
  )

  fill = shape.fill
  fill.solid()
  fill.fore_color.rgb = RGBColor(2, 84, 100)


  left = Inches(transform_inches(10.75 + 0.65 + 0.65 + 6.11 + 0.65))
  top = Inches(transform_inches(0.65))
  width = Inches(transform_inches(6.11))
  height = Inches(transform_inches(2.11))

  shape = slide.shapes.add_shape(
      MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

  text_frame = shape.text_frame
  text_frame.text = f"{number_reviewers(df)}"
  text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

  paragraph = text_frame.add_paragraph()
  paragraph.text = "active code reviewers"
  paragraph.alignment = PP_ALIGN.CENTER
  run = paragraph.runs[0]
  run.font.size = Pt(7)

  fill = shape.fill
  fill.solid()
  fill.fore_color.rgb = RGBColor(2, 84, 100)

  #plots
  left = Inches(transform_inches(0.65))
  top = Inches(transform_inches(0.65 + 0.65 + 1.08 + 0.9))
  width = Inches(transform_inches(10.75))
  height = Inches(transform_inches(8.03))

  pic = slide.shapes.add_picture('./bar_plot.png', left, top, width=width, height=height)

  left = Inches(transform_inches(0.65 + 0.65 + 10.75))
  top_table = (8.03 - (1 + (0.5 * filas_tabla))) / 2
  top = Inches(transform_inches(0.65 + 0.65 + 1.08 + 0.9 + top_table))
  
  height = Inches(transform_inches(1 + 0.5 * filas_tabla))
  width = Inches(transform_inches(12.8))
  pic = slide.shapes.add_picture('./table_language.png', left, top, width=width, height=height)

  left = Inches(transform_inches(0.65))
  top = Inches(transform_inches(0.65 + 0.65 + 1.78 + 0.9 + 8.03))
  height = Inches(transform_inches(5.68))
  width = Inches(transform_inches(24.14))

  pic = slide.shapes.add_picture('./heatmap.png', left, top, width=width, height=height)

  binary_output = BytesIO()
  prs.save(binary_output)
  return binary_output 

cmap, cmap2 = create_cmap("#025464", "#E57C23", "#E8AA42", "#F8F1F1")